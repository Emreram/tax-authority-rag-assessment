"""Build assessment_AI_USE_emresemerci_v2.pptx — full redo, 19 slides.

Single unified visual design across all slides. Narrative voice focused on
*why* I made each dev-process choice with AI. Short bullets, no walls of text.

Design system:
  - Dark navy background, brand-orange left accent bar
  - Eyebrow (orange uppercase 11pt) + Title (white 30pt) + Subtitle (gray italic 15pt)
  - Bullets: 14pt white, max 12 words each
  - "Why I did this" callout box on slides where the reason is the point
  - Footer: "N of 19 · Enterprise RAG Architecture · Emre Semerci"

Run:
    python build_updated_deck.py
Output:
    assessment_AI_USE_emresemerci_v2.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
OUTPUT = HERE / "assessment_AI_USE_emresemerci_v2.pptx"

# ─── Brand palette ───
BD_ORANGE = RGBColor(0xE1, 0x70, 0x00)
BD_ORANGE_DIM = RGBColor(0xB8, 0x5A, 0x00)
BD_NAVY = RGBColor(0x01, 0x68, 0x9B)
BG_DARK = RGBColor(0x0B, 0x12, 0x20)
SURFACE = RGBColor(0x14, 0x1B, 0x2C)
INK = RGBColor(0xE5, 0xE7, 0xEB)
TEXT_DIM = RGBColor(0xB6, 0xC1, 0xCE)
TEXT_MUTED = RGBColor(0x8A, 0x95, 0xA5)

TOTAL = 19


# ═══════════════════════════════════════════════════════════════════════════
# Building blocks (consistent across every slide)
# ═══════════════════════════════════════════════════════════════════════════

def new_slide(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG_DARK; bg.line.fill.background()
    # Accent bar left
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), prs.slide_height)
    bar.fill.solid(); bar.fill.fore_color.rgb = BD_ORANGE; bar.line.fill.background()
    return slide


def add_text(slide, text, *, left, top, width, height, size, color,
             bold=False, italic=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return box


def add_eyebrow(slide, text):
    add_text(slide, text.upper(), left=0.55, top=0.45, width=12, height=0.35,
             size=11, color=BD_ORANGE, bold=True)


def add_title(slide, text, *, top=0.85, size=30):
    add_text(slide, text, left=0.55, top=top, width=12.3, height=0.95,
             size=size, color=INK, bold=True)


def add_subtitle(slide, text, *, top=1.65):
    add_text(slide, text, left=0.55, top=top, width=12.3, height=0.5,
             size=15, color=TEXT_DIM, italic=True)


def add_bullets(slide, bullets, *, left=0.55, top=2.45, width=12.3, height=4.0, size=14, gap=8):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        run = p.add_run()
        run.text = "•   " + b
        run.font.size = Pt(size)
        run.font.color.rgb = INK
        run.font.name = "Calibri"


def add_two_columns(slide, head_l, bullets_l, head_r, bullets_r, *, top=2.45):
    add_column(slide, left=0.55, top=top, width=6.05, heading=head_l, bullets=bullets_l)
    add_column(slide, left=6.85, top=top, width=6.05, heading=head_r, bullets=bullets_r)


def add_column(slide, *, left, top, width, heading, bullets, height=4.0, body_size=13):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p_head = tf.paragraphs[0]
    head_run = p_head.add_run()
    head_run.text = heading
    head_run.font.size = Pt(13)
    head_run.font.bold = True
    head_run.font.color.rgb = BD_ORANGE
    head_run.font.name = "Calibri"
    for b in bullets:
        p = tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = "•   " + b
        run.font.size = Pt(body_size)
        run.font.color.rgb = INK
        run.font.name = "Calibri"


def add_why_box(slide, text, *, top=5.65, height=1.25):
    """Orange-bordered callout: 'Why I did this:' explanation."""
    rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.55), Inches(top), Inches(12.3), Inches(height),
    )
    rect.fill.solid(); rect.fill.fore_color.rgb = SURFACE
    rect.line.color.rgb = BD_ORANGE
    rect.line.width = Pt(1)
    # Inset textbox
    tb = slide.shapes.add_textbox(
        Inches(0.85), Inches(top + 0.15),
        Inches(11.7), Inches(height - 0.30),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    # Label
    p = tf.paragraphs[0]
    label_run = p.add_run()
    label_run.text = "WHY I DID THIS  "
    label_run.font.size = Pt(10)
    label_run.font.bold = True
    label_run.font.color.rgb = BD_ORANGE
    label_run.font.name = "Calibri"
    body_run = p.add_run()
    body_run.text = text
    body_run.font.size = Pt(13)
    body_run.font.italic = True
    body_run.font.color.rgb = INK
    body_run.font.name = "Calibri"


def add_footer(slide, n):
    add_text(slide,
             f"{n} of {TOTAL}    ·    Enterprise RAG Architecture    ·    Emre Semerci",
             left=0.55, top=7.10, width=12.3, height=0.3,
             size=9.5, color=TEXT_MUTED, italic=False)


# ═══════════════════════════════════════════════════════════════════════════
# Slides
# ═══════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    slide = new_slide(prs)
    # Centered title block
    add_text(slide, "ASSESSMENT PRESENTATION",
             left=0.55, top=2.15, width=12.3, height=0.5,
             size=12, color=BD_ORANGE, bold=True)
    add_text(slide, "Enterprise RAG Architecture",
             left=0.55, top=2.65, width=12.3, height=1.0,
             size=46, color=INK, bold=True)
    add_text(slide, "Dutch National Tax Authority",
             left=0.55, top=3.55, width=12.3, height=0.6,
             size=22, color=TEXT_DIM)
    add_text(slide, "How a multi-agent AI workflow turned an architecture brief\ninto a running, hardened demo.",
             left=0.55, top=4.50, width=12.3, height=1.0,
             size=15, color=TEXT_DIM, italic=True)
    # Author block
    add_text(slide, "Emre Semerci",
             left=0.55, top=6.10, width=12.3, height=0.4,
             size=14, color=INK, bold=True)
    add_text(slide, "Lead AI Engineer Technical Assessment   ·   April–May 2026",
             left=0.55, top=6.45, width=12.3, height=0.4,
             size=12, color=TEXT_MUTED)
    add_footer(slide, 1)


def slide_02_brief(prs):
    slide = new_slide(prs)
    add_eyebrow(slide, "The brief")
    add_title(slide, "What the assessment really asks")
    add_subtitle(slide, "Four modules — but the grading lives between the lines.")
    add_two_columns(
        slide,
        "On the page",
        [
            "Module 1 — Ingestion at scale, hierarchical chunking",
            "Module 2 — Hybrid retrieval, reranking, top-K",
            "Module 3 — Agentic CRAG, fail-soft state machine",
            "Module 4 — Cache, RBAC, eval gates",
        ],
        "Between the lines",
        [
            "Specificity over generality — exact numbers, not options",
            "Security first — RBAC must be mathematically leak-proof",
            "Production awareness — TTFT, OOM, quantization math",
            "Domain understanding — Dutch legal hierarchy, ECLI references",
        ],
    )
    add_why_box(slide,
                "I read the brief twice before writing anything. The hidden criteria — "
                "specificity, security math, production constraints — decide the score "
                "more than the visible ones do.")
    add_footer(slide, 2)


def slide_03_workflow(prs):
    slide = new_slide(prs)
    add_eyebrow(slide, "My approach")
    add_title(slide, "Six AI-assisted stages, each with a different tool")
    add_subtitle(slide, "No single agent saw the whole picture uncritically.")
    # Stage boxes in a horizontal flow
    stages = [
        ("1", "Decode", "ChatGPT 5.4"),
        ("2", "Plan", "Claude Code"),
        ("3", "Build", "Claude Code"),
        ("4", "Review", "Hermes Loop"),
        ("5", "Harden", "Claude Code"),
        ("6", "Implement", "Claude Code"),
    ]
    box_w = 1.95
    gap = 0.10
    base_left = 0.55
    base_top = 2.55
    box_h = 1.55
    for i, (num, name, tool) in enumerate(stages):
        left = base_left + i * (box_w + gap)
        rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(base_top), Inches(box_w), Inches(box_h),
        )
        rect.fill.solid(); rect.fill.fore_color.rgb = SURFACE
        rect.line.color.rgb = BD_ORANGE
        rect.line.width = Pt(1)
        # Number
        add_text(slide, num, left=left + 0.18, top=base_top + 0.1, width=0.7, height=0.5,
                 size=20, color=BD_ORANGE, bold=True)
        # Name
        add_text(slide, name, left=left + 0.18, top=base_top + 0.55, width=box_w - 0.36, height=0.4,
                 size=15, color=INK, bold=True)
        # Tool
        add_text(slide, tool, left=left + 0.18, top=base_top + 0.95, width=box_w - 0.36, height=0.4,
                 size=11, color=TEXT_MUTED)
    add_why_box(slide,
                "Each stage needs a different AI strength. ChatGPT is strong at deep analysis; "
                "Claude Code at structured execution; an external Hermes loop at independent review. "
                "Using the same model for everything is how subtle bugs survive review.",
                top=4.35, height=1.55)
    add_footer(slide, 3)


def slide_04_decode(prs):
    slide = new_slide(prs)
    add_eyebrow(slide, "Stage 1 · Decode")
    add_title(slide, "Understand the brief before writing a line of code")
    add_subtitle(slide, "ChatGPT 5.4 with extended thinking, one focused prompt.")
    add_bullets(slide, [
        "Asked: \"give me a master plan and tell me what the brief really grades.\"",
        "Got back: hidden criteria, difficulty per module, 15 locked decisions.",
        "Surprise: Module 3 (CRAG) carries ~35% of the assessor's impression.",
        "Surprise: the RBAC question expects a probability proof, not a policy.",
        "Surprise: the cache threshold tests fiscal-data risk awareness.",
    ])
    add_why_box(slide,
                "If I had started building, I would have spent my budget on Module 1 (the easiest). "
                "Decoding first told me where the marginal hour buys the most score.")
    add_footer(slide, 4)


def slide_05_plan(prs):
    slide = new_slide(prs)
    add_eyebrow(slide, "Stage 2 · Plan")
    add_title(slide, "Plan mode — commit to structure before prose")
    add_subtitle(slide, "Claude Code in plan mode turned the master plan into a work order.")
    add_bullets(slide, [
        "Output: 12-section plan, repo skeleton, 25 file paths, build order.",
        "Locked 15 decisions: HNSW m=16, RRF k=60, cosine ≥ 0.97, e5-large.",
        "Build order: hardest first — schemas, then pseudocode, then prose.",
        "No code yet. Just structure and dependencies.",
    ])
    add_why_box(slide,
                "Without a plan, AI writes the easy stuff first and gets stuck on the hard parts. "
                "With a plan, the schemas exist before any draft references them — "
                "so every reference resolves to a real file. Zero placeholder TODOs.")
    add_footer(slide, 5)


def slide_06_build(prs):
    slide = new_slide(prs)
    add_eyebrow(slide, "Stage 3 · Build")
    add_title(slide, "Sub-agents per artifact type, cross-checked at every step")
    add_subtitle(slide, "Each artifact type got a focused, idempotent build pass.")
    add_two_columns(
        slide,
        "What got built",
        [
            "3 schemas: chunk metadata, OpenSearch mapping, RBAC roles",
            "5 pseudocode files: ingestion, retrieval, CRAG, grader, cache",
            "4 architecture diagrams in mermaid + prose",
            "4 prompt templates: grader, generator, HyDE, decomposition",
        ],
        "Numbers",
        [
            "4,233 lines of pseudocode",
            "1,398 lines of diagrams",
            "24+ files in the final repo",
            "5 explicit build steps, 18 stated assumptions",
        ],
    )
    add_why_box(slide,
                "Sub-agents kept each pass scoped. Cross-checks after every step caught the "
                "inconsistencies AI normally produces (field names that drift, function "
                "signatures that don't match across files).")
    add_footer(slide, 6)


def slide_07_review(prs):
    slide = new_slide(prs)
    add_eyebrow(slide, "Stage 4 · Review")
    add_title(slide, "Two reviewers, two different blind spots caught")
    add_subtitle(slide, "ChatGPT scored quantitatively. Hermes reviewed for tone and clarity.")
    add_two_columns(
        slide,
        "ChatGPT 5.4 — strict scorer",
        [
            "Verdict: 86% complete, 8/10",
            "Found: 14 vs 22 metadata fields contradiction",
            "Found: SQ8 vs fp16 sizing inconsistency",
            "Found: TTFT retry math contradiction",
        ],
        "Hermes Agent — teacher persona",
        [
            "Verdict: \"strong technically, medium on clarity\"",
            "Praised: real decisions, security reasoning, domain awareness",
            "Flagged: 10.7K words causing reviewer fatigue",
            "Flagged: performative tone (\"crown jewel\", \"most candidates\")",
        ],
    )
    add_why_box(slide,
                "I cannot review my own work objectively. Two independent reviewers with different "
                "personas caught two different sets of issues — together they produced the "
                "complete improvement list that became Stage 5.")
    add_footer(slide, 7)


def slide_08_hermes(prs):
    slide = new_slide(prs)
    add_eyebrow(slide, "Stage 4 · Hermes detail")
    add_title(slide, "A two-agent loop, on a 15-minute heartbeat")
    add_subtitle(slide, "Continuous review beats one-shot review for catching slow-burn issues.")
    add_bullets(slide, [
        "Main agent (Codex 5.4 high reasoning): plays the teacher and maintains a feedback document.",
        "Sub-agent 1: writes down questions and imperfections about the feedback document.",
        "Sub-agent 2: answers the questions and fixes the imperfections sub-agent 1 found.",
        "Cycle repeats every 15 minutes for several hours, surfacing issues a single pass misses.",
        "Reasoning skills installed up-front to keep all three agents grounded.",
    ])
    add_why_box(slide,
                "A single review captures what's wrong now. A heartbeat loop captures what becomes "
                "wrong as you keep editing. The teacher persona kept the feedback patient and "
                "constructive — usable instead of demoralising.")
    add_footer(slide, 8)


def slide_09_harden(prs):
    slide = new_slide(prs)
    add_eyebrow(slide, "Stage 5 · Harden")
    add_title(slide, "Address every feedback item — contradictions first")
    add_subtitle(slide, "Internal contradictions hurt credibility more than missing features.")
    add_two_columns(
        slide,
        "What got fixed",
        [
            "Metadata fields: 14 → 22, matched the actual schema",
            "Memory sizing: settled on fp16 ~61 GB primary, SQ8 ~31 GB fallback",
            "Retry TTFT: corrected to ~2030 ms, marked as rare path",
            "RRF made the active OpenSearch pipeline config",
        ],
        "What got cut",
        [
            "Submission rewritten: 10,700 → 7,206 words",
            "Removed all performative language",
            "Hallucination claim qualified to \"fail-closed\"",
            "Repository reorganised into 4 navigable zones",
        ],
    )
    add_why_box(slide,
                "Quality is not added in a final pass — it is recovered. After Stage 4 I knew "
                "exactly where the cracks were. Stage 5 was disciplined repair, not creative work.")
    add_footer(slide, 9)


def slide_10_implement_intro(prs):
    slide = new_slide(prs)
    add_eyebrow(slide, "Stage 6 · Implement")
    add_title(slide, "From paper architecture to running demo")
    add_subtitle(slide, "Everything below this point goes beyond the brief.")
    add_two_columns(
        slide,
        "The stack on a laptop",
        [
            "Docker Model Runner with ai/gemma4:E2B (1.5 GB local LLM)",
            "OpenSearch 2.15 single-node + Redis Stack",
            "intfloat/multilingual-e5-small embeddings",
            "FastAPI + SSE streaming, vanilla-JS frontend",
            "No API keys, no network calls at runtime",
        ],
        "What this proves",
        [
            "Every paper claim is now backed by a curl command",
            "TTFT measured live — not assumed",
            "RBAC leak-proof argument provable by trying to break it",
            "Same architecture as the paper, downsized to a laptop",
        ],
    )
    add_why_box(slide,
                "A paper architecture is unverifiable. A running stack is testable. "
                "Building the demo also surfaced 20 failure modes the paper had glossed over.")
    add_footer(slide, 10)


def slide_11_reliability_s1(prs):
    slide = new_slide(prs)
    add_eyebrow(slide, "Reliability · sprint 1")
    add_title(slide, "Stop the bleeding — three fail-OPEN bugs")
    add_subtitle(slide, "Found while building the demo. None would have shown up in a paper review.")
    add_bullets(slide, [
        "Grader on JSON parse error: defaulted to RELEVANT — silent hallucination bypass.",
        "Citation validator: had a graded[:2] fallback — would publish unverified citations.",
        "/health: returned ready=true even when the LLM was unreachable — splash lied.",
        "All three are now fail-CLOSED. Pre-retrieval RBAC audit confirmed leak-proof.",
    ])
    add_why_box(slide,
                "Production engineers think failure-mode first. These three bugs would have "
                "made the demo look fine until the assessor probed — at which point the "
                "zero-hallucination claim would have collapsed in front of them.")
    add_footer(slide, 11)


def slide_12_reliability_s2_5(prs):
    slide = new_slide(prs)
    add_eyebrow(slide, "Reliability · sprints 2–5")
    add_title(slide, "Twenty failure modes audited, then shipped fixes")
    add_subtitle(slide, "Validation, retry, observability, frontend polish, chaos endpoints.")
    add_two_columns(
        slide,
        "Backend (sprints 2–3)",
        [
            "Per-call LLM timeouts: 15 s classify / 60 s generate",
            "Embedder lock + dedicated thread pool — fixes race condition",
            "Cache, memory, audit fail-soft on Redis hiccup",
            "Per-chunk try/except in ingestion — partial-success reporting",
        ],
        "Frontend + ops (sprints 4–5)",
        [
            "Categorised error events { category, message, request_id }",
            "/readyz polling banner — degraded mode visible to user",
            "Reliability counters surfaced on the Quality page",
            "Chaos endpoints to trip the breaker live during the demo",
        ],
    )
    add_why_box(slide,
                "I picked which sprints to ship by demo-risk × assessment-risk. Anything that "
                "could surface as a 500-error during the live walkthrough went first.")
    add_footer(slide, 12)


def slide_13_refuse_classify(prs):
    slide = new_slide(prs)
    add_eyebrow(slide, "Refuse classification")
    add_title(slide, "Three different \"I don't know\" answers")
    add_subtitle(slide, "A refuse without a reason wastes the user's time.")
    # Three colored boxes (compact)
    boxes = [
        ("CORPUS_GAP",
         "Topic isn't in the corpus.\nSignal: top-1 score < 0.80.\nAdvice: ask the corpus owner.",
         RGBColor(0xFE, 0xE2, 0xE2), RGBColor(0x99, 0x1B, 0x1B), RGBColor(0xDC, 0x26, 0x26)),
        ("TIER_GAP",
         "Content exists, but in a higher tier.\nSignal: best match above your tier.\nAdvice: ask access or a colleague.",
         RGBColor(0xFE, 0xD7, 0xAA), RGBColor(0x9A, 0x34, 0x12), RGBColor(0xEA, 0x58, 0x0C)),
        ("SEMANTIC_MISMATCH",
         "Plausible matches, grader unsure.\nSignal: ≥0.78 in tier, 0 graded RELEVANT.\nAdvice: rephrase with article ref.",
         RGBColor(0xDB, 0xEA, 0xFE), RGBColor(0x1E, 0x40, 0xAF), RGBColor(0x25, 0x63, 0xEB)),
    ]
    box_w = 4.05
    gap = 0.20
    for i, (name, body, bg, fg, border) in enumerate(boxes):
        left = 0.55 + i * (box_w + gap)
        rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(2.45), Inches(box_w), Inches(2.85),
        )
        rect.fill.solid(); rect.fill.fore_color.rgb = bg
        rect.line.color.rgb = border; rect.line.width = Pt(1.5)
        tb = slide.shapes.add_textbox(Inches(left + 0.20), Inches(2.60),
                                      Inches(box_w - 0.40), Inches(2.55))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = 0; tf.margin_right = 0
        p = tf.paragraphs[0]
        head_run = p.add_run()
        head_run.text = name
        head_run.font.size = Pt(15); head_run.font.bold = True
        head_run.font.color.rgb = fg; head_run.font.name = "Calibri"
        for line in body.split("\n"):
            para = tf.add_paragraph()
            para.space_before = Pt(4)
            r = para.add_run()
            r.text = line
            r.font.size = Pt(11); r.font.color.rgb = fg; r.font.name = "Calibri"
    add_why_box(slide,
                "Telling someone \"I don't know\" without telling them why is the worst kind of UX. "
                "Each category gets its own message and audit-log entry — so the operator can "
                "act on the gap, not just notice it.",
                top=5.50, height=1.40)
    add_footer(slide, 13)


def slide_14_false_refuses(prs):
    slide = new_slide(prs)
    add_eyebrow(slide, "Calibration")
    add_title(slide, "Fewer false refuses, no weakened safety gates")
    add_subtitle(slide, "Senior judgment is calibrated, not strict-or-lenient.")
    add_two_columns(
        slide,
        "Four orthogonal knobs",
        [
            "Wider candidate pool: top_k 6 → 10",
            "Retry on IRRELEVANT, not only AMBIGUOUS",
            "AMBIGUOUS-promotion last-chance with strict citations",
            "Force HyDE on the retry pass",
        ],
        "What stays untouched",
        [
            "Zero-hallucination tolerance (citation validator)",
            "Pre-retrieval RBAC filter",
            "Fail-CLOSED grader on parse error",
            "Refuse paths still terminate cleanly",
        ],
    )
    add_why_box(slide,
                "It is tempting to fix \"too many refuses\" by lowering the bar. I refused to. "
                "The four knobs widen the path to the gates, but the gates themselves stay shut.")
    add_footer(slide, 14)


def slide_15_corpus(prs):
    slide = new_slide(prs)
    add_eyebrow(slide, "Corpus")
    add_title(slide, "From 12 to 29 documents — real Dutch tax law")
    add_subtitle(slide, "A 12-doc corpus felt like theatre. So I doubled it down with sources.")
    add_two_columns(
        slide,
        "What was added (18 PUBLIC docs)",
        [
            "Wet IB 2001 — art 3.14, hfd 6, MKB-vrijstelling, eigen woning",
            "Wet OB 1968 — tarieven, KOR, vrijstellingen",
            "Successiewet — vrijstellingen, tarieven",
            "AWR navordering, Awb bezwaar termijn, WKR, WOZ",
            "Two ECLI rulings re-ingested with full structure",
        ],
        "Why these documents",
        [
            "Cover topics that real users would actually ask about",
            "Each chunk carries parent_chunk_id + hierarchy_path",
            "Citations resolve to the exact Lid",
            "The original failing query now answers correctly",
        ],
    )
    add_why_box(slide,
                "The original failing demo question — \"Wanneer mag ik geen huishoudelijke "
                "uitgaven aftrekken?\" — was failing because art 3.14 was simply not in the corpus. "
                "Adding it (and 17 others) made the demo honest.")
    add_footer(slide, 15)


def slide_16_tour(prs):
    slide = new_slide(prs)
    add_eyebrow(slide, "Onboarding")
    add_title(slide, "Guided tour — 10 stops, 4 modules, 75 seconds")
    add_subtitle(slide, "A UI without context is a wall. One button breaks it.")
    add_bullets(slide, [
        "Single button \"Take the tour\" lives on the Quality workspace.",
        "Spotlight + tooltip overlay walks through every page in order.",
        "Each stop names which assessment module it answers.",
        "Vanilla JS, native CSS keyframes, no React, no video framework.",
        "Esc exits, ← / → navigate, localStorage remembers completion.",
    ])
    add_why_box(slide,
                "I stripped the UI of explanatory text deliberately — production apps don't "
                "talk to themselves. The tour is the opt-in version of that text, only "
                "shown when the visitor asks for it.")
    add_footer(slide, 16)


def slide_17_architecture(prs):
    slide = new_slide(prs)
    add_eyebrow(slide, "Architecture")
    add_title(slide, "The whole stack in one frame")
    add_subtitle(slide, "Eight components, one math identity that decides RBAC.")
    # Horizontal flow
    flow = ["Query", "API + OIDC", "Cache", "CRAG loop", "Hybrid retrieval", "Cited answer"]
    box_w = 1.80
    gap = 0.20
    base_left = 0.55
    base_top = 2.40
    for i, name in enumerate(flow):
        left = base_left + i * (box_w + gap)
        rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(base_top), Inches(box_w), Inches(0.8),
        )
        rect.fill.solid(); rect.fill.fore_color.rgb = SURFACE
        rect.line.color.rgb = BD_NAVY; rect.line.width = Pt(1)
        add_text(slide, name, left=left, top=base_top + 0.20, width=box_w, height=0.4,
                 size=12, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_two_columns(
        slide,
        "Decisions across the four modules",
        [
            "Structure-aware chunking, 22-field metadata",
            "Hybrid BM25 + kNN, RRF k=60, optional rerank",
            "Self-hosted Mixtral via vLLM (paper) / Gemma4 (demo)",
            "Tier-partitioned semantic cache cosine ≥ 0.97",
        ],
        "The math that justifies pre-retrieval RBAC",
        [
            "If you post-filter top-40 results with 95% miss rate per chunk:",
            "P(no leak) = 0.95⁴⁰ ≈ 0.13",
            "P(at least one leak) ≈ 0.87",
            "Pre-retrieval makes P(leak) = 0 by construction.",
        ],
        top=3.45,
    )
    add_footer(slide, 17)


def slide_18_deliverables(prs):
    slide = new_slide(prs)
    add_eyebrow(slide, "Deliverables")
    add_title(slide, "Five zones — answer, evidence, ops, reference, demo")
    add_subtitle(slide, "Organised so the assessor can navigate at any depth.")
    add_two_columns(
        slide,
        "Paper architecture",
        [
            "Final submission v2 — 7,206 words, 4 modules in depth",
            "5 pseudocode files, 4 architecture diagrams, 4 prompts",
            "22-field metadata schema, OpenSearch mapping, RBAC model",
            "Metrics matrix, golden test set, CI gate stub",
        ],
        "Runnable demo",
        [
            "Live laptop stack, 29 docs, 174 chunks",
            "Full CRAG pipeline + reliability hardening",
            "Refuse classification with 3 categories",
            "Guided tour, 10 stops covering every module",
        ],
    )
    add_why_box(slide,
                "An assessor opening the README finds one click to the architecture, "
                "one click to the demo, and one click to the bonus performance work. "
                "Every other artifact is one click further. That structure is the deliverable.")
    add_footer(slide, 18)


def slide_19_close(prs):
    slide = new_slide(prs)
    add_eyebrow(slide, "Closing")
    add_title(slide, "What I learned about building with AI agents")
    add_subtitle(slide, "Three things I would tell my past self before starting again.")
    add_bullets(slide, [
        "Different agents have different blind spots — never trust one in a single pass.",
        "Plan structure before prose. Plan mode is the highest-leverage 30 minutes.",
        "Internal contradictions hurt credibility more than missing features do.",
        "A running demo turns architectural answers into engineering handoffs.",
        "The fastest path to quality is structured + reviewed, not bigger model + retry.",
    ])
    add_why_box(slide,
                "The design answer is conservative where the risk is highest, specific where "
                "the brief asks for exact settings, and built through a multi-agent process "
                "that caught its own mistakes before any human reviewer needed to.")
    add_footer(slide, 19)


# ═══════════════════════════════════════════════════════════════════════════
# Main
# ═══════════════════════════════════════════════════════════════════════════

def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    builders = [
        slide_01_title, slide_02_brief, slide_03_workflow, slide_04_decode,
        slide_05_plan, slide_06_build, slide_07_review, slide_08_hermes,
        slide_09_harden, slide_10_implement_intro, slide_11_reliability_s1,
        slide_12_reliability_s2_5, slide_13_refuse_classify, slide_14_false_refuses,
        slide_15_corpus, slide_16_tour, slide_17_architecture, slide_18_deliverables,
        slide_19_close,
    ]
    for b in builders:
        b(prs)
    prs.save(str(OUTPUT))
    print(f"OK — wrote {len(builders)} slides to {OUTPUT.name}")


if __name__ == "__main__":
    main()
