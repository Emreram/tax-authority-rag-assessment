"""Generate operations_justification.pptx from operations_justification.md.

Run from the slides/ directory with python-pptx installed:
    python -m venv .venv && source .venv/Scripts/activate
    pip install -r requirements.txt
    python build_slides.py

Output: slides/output/operations_justification.pptx
"""

from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
SOURCE_MD = HERE / "operations_justification.md"
OUT_DIR = HERE / "output"
OUT_PPTX = OUT_DIR / "operations_justification.pptx"

BD_ORANGE = RGBColor(0xE1, 0x70, 0x00)
BD_NAVY = RGBColor(0x01, 0x68, 0x9B)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x55, 0x5F, 0x6D)
SURFACE = RGBColor(0xF4, 0xF6, 0xF8)


# Field labels: tries each candidate in order. Keeps NL labels working
# for backwards-compatibility with older slide-md revisions.
TITLE_LABELS = ("Title", "Titel")
NOTES_LABELS = ("Speaker notes", "Spreker-notes")
ANCHOR_LABELS = ("UI anchor", "UI-anker")


def parse_slides(md_text: str) -> list[dict]:
    """Split the markdown into slide dicts.

    Recognised sections per slide (EN labels preferred, NL accepted as fallback):
      ## Slide N — <Tab>
      **Title:** ...
      **Bullets (max 3):**
      - ...
      **Speaker notes:**
      <prose>
      **UI anchor:** ...
    """
    slide_blocks = re.split(r"^## Slide \d+ — ", md_text, flags=re.MULTILINE)[1:]
    slides: list[dict] = []
    for block in slide_blocks:
        first_line, rest = block.split("\n", 1)
        tab_name = first_line.strip()

        title = _extract_first_field(rest, TITLE_LABELS)
        bullets = _extract_bullets(rest)
        notes = _extract_first_field(rest, NOTES_LABELS, multiline=True)
        anchor = _extract_first_field(rest, ANCHOR_LABELS)

        slides.append(
            {
                "tab": tab_name,
                "title": title,
                "bullets": bullets,
                "notes": notes,
                "anchor": anchor,
            }
        )
    return slides


def _extract_first_field(text: str, labels: tuple[str, ...], multiline: bool = False) -> str:
    """Try each label until one matches. First non-empty value wins."""
    for label in labels:
        value = _extract_field(text, label, multiline=multiline)
        if value:
            return value
    return ""


def _extract_field(text: str, label: str, multiline: bool = False) -> str:
    pattern = rf"\*\*{re.escape(label)}:\*\*\s*(.*?)(?=\n\s*\n|\s*\Z)"
    flags = re.DOTALL if multiline else 0
    m = re.search(pattern, text, flags=flags)
    if not m:
        return ""
    value = m.group(1).strip()
    # Collapse interior newlines+whitespace into single spaces for non-multiline
    if not multiline:
        value = re.sub(r"\s+", " ", value)
    return value


def _extract_bullets(text: str) -> list[str]:
    m = re.search(r"\*\*Bullets[^*]*\*\*\s*\n((?:- .+\n?)+)", text)
    if not m:
        return []
    raw = m.group(1).strip().splitlines()
    bullets: list[str] = []
    for line in raw:
        if line.startswith("- "):
            bullets.append(line[2:].strip())
        elif bullets:
            bullets[-1] += " " + line.strip()
    return bullets


def build_presentation(slides: list[dict]) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # fully blank

    for s in slides:
        slide = prs.slides.add_slide(blank_layout)
        _add_accent_bar(slide)
        _add_eyebrow(slide, s["tab"])
        _add_title(slide, s["title"])
        _add_bullets(slide, s["bullets"])
        _add_anchor_footer(slide, s["anchor"])
        _add_speaker_notes(slide, s["notes"])

    return prs


def _add_accent_bar(slide) -> None:
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.25), Inches(7.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BD_ORANGE
    bar.line.fill.background()


def _add_eyebrow(slide, tab_name: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(11), Inches(0.4))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"OPERATIONS · {tab_name.upper()}"
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = BD_NAVY


def _add_title(slide, title_text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(11.5), Inches(1.4))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = INK


def _add_bullets(slide, bullets: list[str]) -> None:
    box = slide.shapes.add_textbox(Inches(0.7), Inches(2.7), Inches(11.5), Inches(4.0))
    tf = box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(14)
        label, _, body = bullet.partition(":")
        if body:
            run_label = p.add_run()
            run_label.text = f"{label.strip()}: "
            run_label.font.size = Pt(20)
            run_label.font.bold = True
            run_label.font.color.rgb = BD_ORANGE
            run_body = p.add_run()
            run_body.text = body.strip()
            run_body.font.size = Pt(20)
            run_body.font.color.rgb = INK
        else:
            run = p.add_run()
            run.text = bullet
            run.font.size = Pt(20)
            run.font.color.rgb = INK


def _add_anchor_footer(slide, anchor_text: str) -> None:
    box = slide.shapes.add_textbox(
        Inches(0.7), Inches(6.85), Inches(11.5), Inches(0.4)
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"UI anchor: {anchor_text}"
    run.font.size = Pt(11)
    run.font.italic = True
    run.font.color.rgb = MUTED


def _add_speaker_notes(slide, notes_text: str) -> None:
    if not notes_text:
        return
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = notes_text


def main() -> None:
    if not SOURCE_MD.exists():
        raise SystemExit(f"Source file not found: {SOURCE_MD}")
    md_text = SOURCE_MD.read_text(encoding="utf-8")
    slides = parse_slides(md_text)
    if len(slides) != 5:
        raise SystemExit(
            f"Expected 5 slide sections in {SOURCE_MD.name}, found {len(slides)}."
        )
    prs = build_presentation(slides)
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPTX)
    print(f"OK — {len(slides)} slides written to {OUT_PPTX.relative_to(HERE.parent)}")


if __name__ == "__main__":
    main()
